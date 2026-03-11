from __future__ import annotations

import hashlib
import io
import json
import os
import re
import secrets
import subprocess
import sys
import threading
import tkinter as tk
from pathlib import Path
from tkinter import filedialog, messagebox

import pandas as pd
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Cm, Pt

INVALID_FILENAME_CHARS = r'[\\/:*?"<>|]'
QR_SIZE_CM = 2.0
QR_MARGIN_CM = 0.6


def sanitize_filename(nombre: str, index: int, max_len: int = 120) -> str:
    base = f"{index:02d} _Constancia_Poster_{nombre}".strip()
    base = re.sub(INVALID_FILENAME_CHARS, "", base).strip().strip(".")
    if not base:
        base = f"{index:02d} _Constancia_Poster"
    max_base_len = max_len - len(".pptx")
    if len(base) > max_base_len:
        base = base[:max_base_len].rstrip().rstrip(".")
    return f"{base}.pptx"


def replace_in_paragraph(paragraph, replacements: dict[str, str]) -> bool:
    if not paragraph.runs:
        return False

    original_text = "".join(run.text for run in paragraph.runs)
    updated_text = original_text
    for key, value in replacements.items():
        updated_text = updated_text.replace(key, value)

    if updated_text == original_text:
        return False

    for run in paragraph.runs:
        run_text = run.text
        new_text = run_text
        for key, value in replacements.items():
            new_text = new_text.replace(key, value)
        if new_text != run_text:
            run.text = new_text

    combined = "".join(run.text for run in paragraph.runs)
    if any(key in combined for key in replacements.keys()):
        paragraph.text = updated_text

    return True


def replace_in_text_frame(text_frame, replacements: dict[str, str]) -> bool:
    changed = False
    for paragraph in text_frame.paragraphs:
        if replace_in_paragraph(paragraph, replacements):
            changed = True
    return changed


def iter_shapes(shapes):
    for shape in shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            for subshape in iter_shapes(shape.shapes):
                yield subshape
        else:
            yield shape


def replace_in_slide(slide, replacements: dict[str, str]) -> None:
    for shape in iter_shapes(slide.shapes):
        if shape.has_text_frame:
            replace_in_text_frame(shape.text_frame, replacements)
        if shape.has_table:
            for row in shape.table.rows:
                for cell in row.cells:
                    replace_in_text_frame(cell.text_frame, replacements)


def exportar_pptx_a_pdf(pptx_path: Path, pdf_path: Path) -> None:
    try:
        import comtypes.client
    except ModuleNotFoundError as exc:
        raise ModuleNotFoundError(
            "No se encontro 'comtypes'. Instala con: pip install comtypes"
        ) from exc

    powerpoint = comtypes.client.CreateObject("PowerPoint.Application")
    powerpoint.Visible = 1
    presentation = powerpoint.Presentations.Open(str(pptx_path), WithWindow=False)
    try:
        presentation.SaveAs(str(pdf_path), 32)  # 32 = ppSaveAsPDF
    finally:
        presentation.Close()
        powerpoint.Quit()

def generar_hash(nombre: str, titulo: str, index: int, secreto: str) -> str:
    base = f"{secreto}|{nombre}|{titulo}|{index}"
    return hashlib.sha256(base.encode("utf-8")).hexdigest().upper()[:16]


def generar_qr_bytes(data: str) -> io.BytesIO:
    try:
        import qrcode
    except ModuleNotFoundError as exc:
        raise ModuleNotFoundError(
            "No se encontro 'qrcode'. Instala con: pip install qrcode[pil]"
        ) from exc

    qr = qrcode.QRCode(
        version=None,
        error_correction=qrcode.constants.ERROR_CORRECT_M,
        box_size=10,
        border=2,
    )
    qr.add_data(data)
    qr.make(fit=True)
    img = qr.make_image(fill_color="black", back_color="white")
    buffer = io.BytesIO()
    img.save(buffer, format="PNG")
    buffer.seek(0)
    return buffer


def encontrar_placeholder_qr(slide):
    for shape in iter_shapes(slide.shapes):
        if not shape.has_text_frame:
            continue
        text = shape.text_frame.text or ""
        if "{{QR}}" in text:
            return shape
    return None


def agregar_qr_y_folio(slide, qr_data: str, folio: str, slide_width, slide_height) -> None:
    placeholder = encontrar_placeholder_qr(slide)
    qr_stream = generar_qr_bytes(qr_data)

    target_qr_size = Cm(QR_SIZE_CM)
    text_height = Cm(0.6)
    gap = Cm(0.1)

    if placeholder is not None:
        placeholder.text_frame.text = placeholder.text_frame.text.replace("{{QR}}", "").strip()
        if placeholder.width <= 0 or placeholder.height <= 0:
            placeholder = None
        else:
            slide.shapes.add_picture(
                qr_stream,
                placeholder.left,
                placeholder.top,
                width=placeholder.width,
                height=placeholder.height,
            )
            folio_top = placeholder.top + placeholder.height + gap
            max_folio_top = slide_height - text_height - gap
            if folio_top > max_folio_top:
                folio_top = max_folio_top
            textbox = slide.shapes.add_textbox(
                placeholder.left,
                folio_top,
                placeholder.width,
                text_height,
            )
            p = textbox.text_frame.paragraphs[0]
            p.text = f"FOLIO: {folio}"
            p.font.size = Pt(8)
            p.alignment = PP_ALIGN.CENTER
            return

    qr_size = target_qr_size
    margin = Cm(QR_MARGIN_CM)
    left = margin
    top = slide_height - qr_size - margin - text_height - gap

    slide.shapes.add_picture(qr_stream, left, top, width=qr_size, height=qr_size)

    textbox = slide.shapes.add_textbox(left, top + qr_size + gap, qr_size, text_height)
    p = textbox.text_frame.paragraphs[0]
    p.text = f"FOLIO: {folio}"
    p.font.size = Pt(8)
    p.alignment = PP_ALIGN.CENTER


def generar_constancias(
    csv_path: Path,
    template_path: Path,
    output_dir: Path,
    exportar_pdf: bool,
    incluir_qr: bool,
    url_base: str,
    secreto: str,
) -> tuple[int, int]:
    df = pd.read_csv(csv_path, encoding="utf-8", dtype=str, keep_default_na=False)
    if "NOMBRE" not in df.columns or "TITULO" not in df.columns:
        raise ValueError("El CSV debe contener las columnas NOMBRE y TITULO.")

    output_dir.mkdir(parents=True, exist_ok=True)
    pdf_dir = output_dir / "pdf"
    if exportar_pdf:
        pdf_dir.mkdir(parents=True, exist_ok=True)

    total = 0
    pdf_total = 0
    exportar_pdf_activo = exportar_pdf
    registros_validacion = []

    for index, row in df.iterrows():
        nombre = str(row.get("NOMBRE", "")).strip()
        titulo = str(row.get("TITULO", "")).strip()
        replacements = {
            "{{NOMBRE}}": nombre,
            "{{TITULO}}": titulo,
        }

        presentation = Presentation(template_path)
        for slide in presentation.slides:
            replace_in_slide(slide, replacements)

            if incluir_qr:
                folio = generar_hash(nombre, titulo, index + 1, secreto)
                qr_url = f"{url_base}{folio}"
                agregar_qr_y_folio(slide, qr_url, folio, presentation.slide_width, presentation.slide_height)

        filename = sanitize_filename(nombre, index + 1)
        output_path = output_dir / filename
        presentation.save(output_path)
        total += 1

        pdf_name = None
        if exportar_pdf_activo:
            pdf_path = pdf_dir / Path(filename).with_suffix(".pdf").name
            try:
                exportar_pptx_a_pdf(output_path, pdf_path)
                pdf_total += 1
                pdf_name = pdf_path.name
            except ModuleNotFoundError:
                exportar_pdf_activo = False

        if incluir_qr:
            folio = generar_hash(nombre, titulo, index + 1, secreto)
            registros_validacion.append(
                {
                    "hash": folio,
                    "nombre": nombre,
                    "actividad": titulo,
                    "archivo": pdf_name or output_path.name,
                }
            )

    if incluir_qr:
        hashes_path = output_dir / "hashes_validacion.json"
        with hashes_path.open("w", encoding="utf-8") as f:
            json.dump(registros_validacion, f, ensure_ascii=False, indent=2)

    return total, pdf_total if exportar_pdf else 0


class App(tk.Tk):
    def __init__(self) -> None:
        super().__init__()
        self.title("Generador de Constancias de Posters")
        self.geometry("760x580")
        self.resizable(False, False)

        self.template_var = tk.StringVar()
        self.csv_var = tk.StringVar()
        self.output_var = tk.StringVar()
        self.exportar_pdf_var = tk.BooleanVar(value=True)
        self.incluir_qr_var = tk.BooleanVar(value=True)
        self.url_base_var = tk.StringVar(value="https://semanameca.upiiz.ipn.mx/validar_qr.html?hash=")
        self.secreto_var = tk.StringVar(value="CAMBIA_ESTE_SECRETO")
        self.status_var = tk.StringVar(value="Listo.")
        self.logo_image = None

        self._build_ui()

    def _build_ui(self) -> None:
        self._build_logo()
        title = tk.Label(self, text="Generar constancias de posters", font=("Segoe UI", 14, "bold"))
        title.pack(pady=12)

        frame = tk.Frame(self)
        frame.pack(fill="x", padx=12)

        self._row(frame, "Plantilla PPTX:", self.template_var, self._browse_template, 0)
        self._row(frame, "Archivo CSV:", self.csv_var, self._browse_csv, 1)
        self._row(
            frame,
            "Carpeta de salida:",
            self.output_var,
            self._browse_output,
            2,
            extra_button=("Abrir", self._open_output_dir),
        )

        qr_frame = tk.Frame(frame)
        qr_frame.grid(row=3, column=0, sticky="ew", pady=4)
        qr_frame.columnconfigure(1, weight=1)
        tk.Label(qr_frame, text="URL validación:", width=18, anchor="w").grid(row=0, column=0, padx=(0, 8))
        tk.Entry(qr_frame, textvariable=self.url_base_var).grid(row=0, column=1, sticky="ew")

        secret_frame = tk.Frame(frame)
        secret_frame.grid(row=4, column=0, sticky="ew", pady=4)
        secret_frame.columnconfigure(1, weight=1)
        tk.Label(secret_frame, text="Secreto:", width=18, anchor="w").grid(row=0, column=0, padx=(0, 8))
        tk.Entry(secret_frame, textvariable=self.secreto_var, show="*").grid(row=0, column=1, sticky="ew")
        tk.Button(secret_frame, text="Generar secreto", command=self._sugerir_secreto).grid(
            row=0, column=2, padx=(8, 0)
        )

        options = tk.Frame(self)
        options.pack(fill="x", padx=12, pady=8)
        tk.Checkbutton(options, text="Exportar también a PDF (requiere PowerPoint)", variable=self.exportar_pdf_var).pack(anchor="w")
        tk.Checkbutton(options, text="Incluir QR + folio de seguridad", variable=self.incluir_qr_var).pack(anchor="w")

        actions = tk.Frame(self)
        actions.pack(fill="x", padx=12, pady=10)
        self.run_button = tk.Button(actions, text="Generar", command=self._run)
        self.run_button.pack(side="left")
        self.exit_button = tk.Button(actions, text="Salir", command=self.destroy)
        self.exit_button.pack(side="left", padx=(8, 0))

        status = tk.Label(self, textvariable=self.status_var, anchor="w", fg="#1a1a1a")
        status.pack(fill="x", padx=12, pady=6)

        hint = tk.Label(
            self,
            text=(
                "Nota: la plantilla original no se modifica; se generan archivos nuevos. "
                "Para validar en la web, guarda en Constancias/Posters/salida_pptx."
            ),
            fg="#555555",
            anchor="w",
        )
        hint.pack(fill="x", padx=12)

    def _build_logo(self) -> None:
        base_dir = Path(__file__).resolve().parent
        logo_path = base_dir / "Docs" / "Logo_SM_Capas (20250615102813).png"
        if not logo_path.exists():
            return

        image = None
        try:
            from PIL import Image, ImageTk

            img = Image.open(logo_path)
            img = img.convert("RGBA")
            img = img.resize((180, 180), Image.LANCZOS)
            image = ImageTk.PhotoImage(img)
        except Exception:
            try:
                image = tk.PhotoImage(file=str(logo_path))
                image = image.subsample(3, 3)
            except Exception:
                image = None

        if image is None:
            return

        self.logo_image = image
        logo_label = tk.Label(self, image=self.logo_image)
        logo_label.pack(pady=(12, 4))

    def _row(
        self,
        parent,
        label,
        variable,
        command,
        row_idx: int,
        extra_button: tuple[str, callable] | None = None,
    ) -> None:
        row = tk.Frame(parent)
        row.grid(row=row_idx, column=0, sticky="ew", pady=4)
        row.columnconfigure(1, weight=1)

        tk.Label(row, text=label, width=18, anchor="w").grid(row=0, column=0, padx=(0, 8))
        entry = tk.Entry(row, textvariable=variable)
        entry.grid(row=0, column=1, sticky="ew")
        tk.Button(row, text="Buscar...", command=command).grid(row=0, column=2, padx=(8, 0))
        if extra_button:
            text, cmd = extra_button
            tk.Button(row, text=text, command=cmd).grid(row=0, column=3, padx=(8, 0))

    def _browse_template(self) -> None:
        path = filedialog.askopenfilename(
            title="Seleccionar plantilla PPTX",
            filetypes=[("PowerPoint", "*.pptx")],
        )
        if path:
            self.template_var.set(path)

    def _browse_csv(self) -> None:
        path = filedialog.askopenfilename(
            title="Seleccionar archivo CSV",
            filetypes=[("CSV", "*.csv")],
        )
        if path:
            self.csv_var.set(path)

    def _browse_output(self) -> None:
        path = filedialog.askdirectory(title="Seleccionar carpeta de salida")
        if path:
            self.output_var.set(path)

    def _open_output_dir(self) -> None:
        path_str = self.output_var.get().strip()
        if not path_str:
            messagebox.showinfo("Carpeta de salida", "Selecciona una carpeta de salida primero.")
            return

        path = Path(path_str)
        if not path.exists():
            messagebox.showerror("Error", "La carpeta de salida no existe.")
            return

        try:
            if sys.platform.startswith("win"):
                os.startfile(path)  # type: ignore[attr-defined]
            elif sys.platform == "darwin":
                subprocess.run(["open", str(path)], check=False)
            else:
                subprocess.run(["xdg-open", str(path)], check=False)
        except Exception as exc:
            messagebox.showerror("Error", f"No se pudo abrir la carpeta:\n{exc}")

    def _run(self) -> None:
        template_path = Path(self.template_var.get().strip())
        csv_path = Path(self.csv_var.get().strip())
        output_dir = Path(self.output_var.get().strip())
        exportar_pdf = self.exportar_pdf_var.get()
        incluir_qr = self.incluir_qr_var.get()
        url_base = self.url_base_var.get().strip()
        secreto = self.secreto_var.get().strip()

        if not template_path.exists():
            messagebox.showerror("Error", "Selecciona una plantilla PPTX válida.")
            return
        if not csv_path.exists():
            messagebox.showerror("Error", "Selecciona un CSV válido.")
            return
        if not output_dir.exists():
            messagebox.showerror("Error", "Selecciona una carpeta de salida válida.")
            return
        if incluir_qr and not url_base:
            messagebox.showerror("Error", "Ingresa la URL base para validación.")
            return
        if incluir_qr and (not secreto or secreto == "CAMBIA_ESTE_SECRETO"):
            self._sugerir_secreto()
            return

        self.run_button.config(state="disabled")
        self.exit_button.config(state="disabled")
        self.status_var.set("Procesando... Por favor espera.")

        thread = threading.Thread(
            target=self._run_task,
            args=(csv_path, template_path, output_dir, exportar_pdf, incluir_qr, url_base, secreto),
            daemon=True,
        )
        thread.start()

    def _run_task(
        self,
        csv_path: Path,
        template_path: Path,
        output_dir: Path,
        exportar_pdf: bool,
        incluir_qr: bool,
        url_base: str,
        secreto: str,
    ) -> None:
        try:
            total, pdf_total = generar_constancias(
                csv_path,
                template_path,
                output_dir,
                exportar_pdf,
                incluir_qr,
                url_base,
                secreto,
            )
            msg = f"Generadas {total} constancias en: {output_dir}"
            if exportar_pdf:
                msg += f"\nPDFs generados: {pdf_total} en {output_dir / 'pdf'}"
            if incluir_qr:
                msg += f"\nHashes guardados en: {output_dir / 'hashes_validacion.json'}"
            self._finish_success(msg)
        except Exception as exc:
            self._finish_error(str(exc))

    def _finish_success(self, message: str) -> None:
        self.after(0, lambda: self._finish(message, success=True))

    def _finish_error(self, message: str) -> None:
        self.after(0, lambda: self._finish(message, success=False))

    def _finish(self, message: str, success: bool) -> None:
        self.run_button.config(state="normal")
        self.exit_button.config(state="normal")
        self.status_var.set("Listo." if success else "Error.")
        if success:
            messagebox.showinfo("Completado", message)
        else:
            messagebox.showerror("Error", message)

    def _sugerir_secreto(self) -> None:
        sugerido = f"Meca2025-Posters-{secrets.token_urlsafe(10)}"
        self.secreto_var.set(sugerido)
        self.clipboard_clear()
        self.clipboard_append(sugerido)
        messagebox.showinfo(
            "Secreto sugerido",
            "Se genero un secreto y se copio al portapapeles.\n\n"
            f"Secreto:\n{sugerido}\n\n"
            "Guardalo en un lugar seguro. Si lo cambias, los QR anteriores no validaran.",
        )


if __name__ == "__main__":
    app = App()
    app.mainloop()
