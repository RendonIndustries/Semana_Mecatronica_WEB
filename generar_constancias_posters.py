from __future__ import annotations

import re
import sys
from pathlib import Path

import pandas as pd
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

INVALID_FILENAME_CHARS = r'[\\/:*?"<>|]'


def sanitize_filename(nombre: str, index: int, max_len: int = 120) -> str:
    base = f"{index:02d} _Constancia_Poster_{nombre}".strip()
    base = re.sub(INVALID_FILENAME_CHARS, "", base).strip().strip(".")
    if not base:
        base = f"{index:02d} _Constancia_Poster"
    max_base_len = max_len - len(".pptx")
    if len(base) > max_base_len:
        base = base[:max_base_len].rstrip().rstrip(".")
    return f"{base}.pptx"


def replace_in_paragraph(paragraph, replacements: dict[str, str]) -> bool:
    if not paragraph.runs:
        return False

    original_text = "".join(run.text for run in paragraph.runs)
    updated_text = original_text
    for key, value in replacements.items():
        updated_text = updated_text.replace(key, value)

    if updated_text == original_text:
        return False

    for run in paragraph.runs:
        run_text = run.text
        new_text = run_text
        for key, value in replacements.items():
            new_text = new_text.replace(key, value)
        if new_text != run_text:
            run.text = new_text

    combined = "".join(run.text for run in paragraph.runs)
    if any(key in combined for key in replacements.keys()):
        paragraph.text = updated_text

    return True


def replace_in_text_frame(text_frame, replacements: dict[str, str]) -> bool:
    changed = False
    for paragraph in text_frame.paragraphs:
        if replace_in_paragraph(paragraph, replacements):
            changed = True
    return changed


def iter_shapes(shapes):
    for shape in shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            for subshape in iter_shapes(shape.shapes):
                yield subshape
        else:
            yield shape


def replace_in_slide(slide, replacements: dict[str, str]) -> None:
    for shape in iter_shapes(slide.shapes):
        if shape.has_text_frame:
            replace_in_text_frame(shape.text_frame, replacements)
        if shape.has_table:
            for row in shape.table.rows:
                for cell in row.cells:
                    replace_in_text_frame(cell.text_frame, replacements)


def main() -> None:
    base_dir = Path(__file__).resolve().parent
    csv_path = base_dir / "posters.csv"
    template_path = base_dir / "plantilla_constancia.pptx"
    output_dir = base_dir / "salida_pptx"

    if not csv_path.exists():
        print(f"No se encontro el archivo CSV: {csv_path}")
        sys.exit(1)
    if not template_path.exists():
        print(f"No se encontro la plantilla PPTX: {template_path}")
        sys.exit(1)

    output_dir.mkdir(parents=True, exist_ok=True)

    df = pd.read_csv(csv_path, encoding="utf-8", dtype=str, keep_default_na=False)
    if "NOMBRE" not in df.columns or "TITULO" not in df.columns:
        print("El CSV debe contener las columnas NOMBRE y TITULO.")
        sys.exit(1)

    total = 0
    for index, row in df.iterrows():
        nombre = str(row.get("NOMBRE", "")).strip()
        titulo = str(row.get("TITULO", "")).strip()
        replacements = {
            "{{NOMBRE}}": nombre,
            "{{TITULO}}": titulo,
        }

        presentation = Presentation(template_path)
        for slide in presentation.slides:
            replace_in_slide(slide, replacements)

        filename = sanitize_filename(nombre, index + 1)
        output_path = output_dir / filename
        presentation.save(output_path)
        total += 1

    print(f"Generadas {total} constancias en: {output_dir.resolve()}")


if __name__ == "__main__":
    main()
